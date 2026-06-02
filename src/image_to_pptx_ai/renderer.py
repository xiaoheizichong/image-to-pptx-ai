from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from image_to_pptx_ai.colors import hex_to_rgb
from image_to_pptx_ai.models import SlideElement, SlideSpec


SLIDE_SIZES = {
    "wide": (13.333, 7.5),
    "standard": (10.0, 7.5),
}


def render_pptx(spec: SlideSpec, output_path: Path, *, source_image: Path | None = None) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    width, height = SLIDE_SIZES[spec.size]
    presentation.slide_width = Inches(width)
    presentation.slide_height = Inches(height)

    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    _set_background(slide, spec.background)

    for element in spec.elements:
        _render_element(slide, element, source_image=source_image)

    presentation.save(output_path)
    return output_path


def _set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _render_element(slide, element: SlideElement, *, source_image: Path | None) -> None:
    if element.kind == "text":
        _add_text(slide, element)
    elif element.kind == "shape":
        _add_shape(slide, element)
    elif element.kind == "image":
        _add_image(slide, element, source_image=source_image)


def _add_text(slide, element: SlideElement) -> None:
    box = slide.shapes.add_textbox(
        Inches(element.x), Inches(element.y), Inches(element.w), Inches(element.h)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = element.text or ""
    _style_run(run, element)

    for bullet in element.bullets:
        bullet_paragraph = frame.add_paragraph()
        bullet_paragraph.text = bullet
        bullet_paragraph.level = 0
        bullet_paragraph.font.size = Pt(max(8, element.font_size - 2))
        bullet_paragraph.font.color.rgb = _rgb(element.color)


def _style_run(run, element: SlideElement) -> None:
    run.font.size = Pt(element.font_size)
    run.font.bold = element.bold
    run.font.italic = element.italic
    run.font.color.rgb = _rgb(element.color)


def _add_shape(slide, element: SlideElement) -> None:
    shape_type = {
        "rect": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "line": MSO_SHAPE.RECTANGLE,
        None: MSO_SHAPE.RECTANGLE,
    }[element.shape]
    shape = slide.shapes.add_shape(
        shape_type, Inches(element.x), Inches(element.y), Inches(element.w), Inches(element.h)
    )

    if element.fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(element.fill)
    else:
        shape.fill.background()

    if element.outline:
        shape.line.color.rgb = _rgb(element.outline)
    else:
        shape.line.fill.background()


def _add_image(slide, element: SlideElement, *, source_image: Path | None) -> None:
    image_path = element.image_path
    if image_path == "__SOURCE_IMAGE__" and source_image is not None:
        image_path = str(source_image)
    if not image_path:
        return
    path = Path(image_path)
    if not path.exists():
        return
    slide.shapes.add_picture(
        str(path), Inches(element.x), Inches(element.y), width=Inches(element.w), height=Inches(element.h)
    )


def _rgb(color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(color))

